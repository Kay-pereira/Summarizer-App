import traceback
from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import IsAuthenticated
from django.contrib.auth.models import User

import fitz  # PyMuPDF
import docx
from pptx import Presentation
import openai
from decouple import config

from .models import Summary
from .serializers import RegisterSerializer, SummarySerializer

openai.api_key = config("OPEN_AI_API_KEY", default=None)


class RegisterView(APIView):
    authentication_classes = []  # allow unauthenticated
    permission_classes = []      # allow unauthenticated

    def post(self, request):
        serializer = RegisterSerializer(data=request.data)
        if serializer.is_valid():
            user = serializer.save()
            return Response(
                {"message": "Account created", "username": user.username},
                status=status.HTTP_201_CREATED,
            )
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


class SummarizeView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = [MultiPartParser]

    def post(self, request, format=None):
        try:
            # Ensure OpenAI key present
            if not openai.api_key:
                return Response({"error": "OpenAI API key not configured."}, status=500)

            file = request.FILES.get("file")
            if not file:
                return Response({"error": "No file uploaded"}, status=400)

            ext = file.name.lower()
            if ext.endswith(".pdf"):
                text = self.extract_pdf(file)
            elif ext.endswith(".docx"):
                text = self.extract_docx(file)
            elif ext.endswith(".pptx"):
                text = self.extract_pptx(file)
            else:
                return Response({"error": "Unsupported file type"}, status=400)

            if len(text.strip()) == 0:
                return Response({"error": "File has no readable text."}, status=400)

            # Limit text length to avoid hitting token limits
            truncated = text[:4000]
            summary = self.get_summary(truncated)

            # Save to DB tied to user
            record = Summary.objects.create(
                user=request.user,
                file_name=file.name,
                original_text=truncated,
                summary_text=summary.get("summary", summary if isinstance(summary, str) else "")
            )

            output = {
                "overview": summary.get("overview") if isinstance(summary, dict) else "",
                "summary": summary.get("summary") if isinstance(summary, dict) else (summary if isinstance(summary, str) else ""),
                "id": record.id,
            }

            return Response(output, status=200)

        except ValueError as ve:
            return Response({"error": str(ve)}, status=400)
        except Exception as e:
            trace = traceback.format_exc()
            print("❌ SummarizeView Exception:", str(e))
            print(trace)
            return Response({"error": "Internal server error", "trace": str(e)}, status=500)

    def extract_pdf(self, file):
        try:
            data = file.read()
            doc = fitz.open(stream=data, filetype="pdf")
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            return "\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Invalid or corrupted PDF file: {e}")

    def extract_docx(self, file):
        try:
            doc = docx.Document(file)
            return "\n".join([p.text for p in doc.paragraphs if p.text])
        except Exception as e:
            raise ValueError(f"Invalid or corrupted DOCX file: {e}")

    def extract_pptx(self, file):
        try:
            prs = Presentation(file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            raise ValueError(f"Invalid or corrupted PPTX file: {e}")

    def get_summary(self, text):
        # Use ChatCompletion (older style) — keep same as in your code; if you use new OpenAI SDK, adapt here.
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful teaching assistant."},
                    {"role": "user", "content": f"Please provide an overview and a detailed summary of the following lesson content:\n\n{text}"}
                ]
            )
            result = response.choices[0].message.content
            return {"overview": result.split("\n")[0] if result else "", "summary": result}
        except Exception as e:
            # bubble up as an error to be returned as 500
            raise Exception(f"OpenAI API error: {e}")


class SummaryListView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        qs = Summary.objects.filter(user=request.user).order_by("-created_at")
        serializer = SummarySerializer(qs, many=True)
        return Response(serializer.data, status=200)
